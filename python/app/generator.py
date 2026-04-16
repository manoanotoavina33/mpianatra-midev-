from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import random

# Couleurs
FOND_SOMBRE = RGBColor(15, 23, 42)
FOND_CLAIR = RGBColor(30, 41, 59)
ORANGE = RGBColor(251, 146, 60)
BLANC = RGBColor(255, 255, 255)
JAUNE = RGBColor(251, 191, 36)
VIOLET = RGBColor(99, 102, 241)
VERT = RGBColor(34, 197, 94)
GRIS = RGBColor(148, 163, 184)

SALUTATIONS = [
    "Bienvenue ! 👋",
    "Bonjour ! 😊",
    "Bonsoir ! 🌙",
    "Salut ! ✨",
    "Hello ! 🚀",
]

def set_fond(slide, couleur):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = couleur

def ajouter_texte(slide, texte, x, y, w, h, taille,
                  couleur, bold=False,
                  align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(texte)
    run.font.size = Pt(taille)
    run.font.bold = bold
    run.font.color.rgb = couleur

def slide_salutation(prs):
    """Slide de bienvenue avec salutation aléatoire"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, FOND_SOMBRE)

    salut = random.choice(SALUTATIONS)

    # Cercle décoratif simulé avec rectangle arrondi
    shape = slide.shapes.add_shape(
        9, Inches(3.5), Inches(1.5), Inches(3), Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    ajouter_texte(
        slide, salut,
        1, 2.5, 8, 1.5,
        48, BLANC, bold=True,
        align=PP_ALIGN.CENTER
    )
    ajouter_texte(
        slide, "Mpianatra-midev",
        1, 4.2, 8, 1,
        20, GRIS,
        align=PP_ALIGN.CENTER
    )

def slide_titre(prs, titre, auteur, institution, annee):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, FOND_SOMBRE)

    # Ligne décorative
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(2.2), Inches(10), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    ajouter_texte(
        slide, titre,
        0.5, 0.4, 9, 1.8,
        32, ORANGE, bold=True,
        align=PP_ALIGN.CENTER
    )

    infos = []
    if auteur:
        infos.append(f"👤 {auteur}")
    if institution:
        infos.append(f"🏫 {institution}")
    if annee:
        infos.append(f"📅 {annee}")

    y = 2.5
    for info in infos:
        ajouter_texte(
            slide, info,
            1, y, 8, 0.6,
            16, BLANC,
            align=PP_ALIGN.CENTER
        )
        y += 0.6

    ajouter_texte(
        slide, "Généré par Mpianatra-midev",
        1, 6.5, 8, 0.6,
        11, GRIS,
        align=PP_ALIGN.CENTER
    )

def slide_sommaire(prs, slides_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, FOND_SOMBRE)

    ajouter_texte(
        slide, "📋 Sommaire",
        0.5, 0.2, 9, 0.9,
        28, ORANGE, bold=True
    )

    # Ligne séparatrice
    shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    y = 1.3
    num = 1
    for s in slides_data:
        if s.get("type") not in ["titre", "sommaire", "remerciements", "salutation"]:
            ajouter_texte(
                slide,
                f"{num}.  {s.get('titre', '')}",
                0.8, y, 8.5, 0.55,
                15, BLANC
            )
            y += 0.58
            num += 1
            if num > 9:
                break

def slide_contenu(prs, titre, contenu, points=None, couleur_accent=None):
    if couleur_accent is None:
        couleur_accent = ORANGE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, FOND_SOMBRE)

    # Numéro de slide décoratif
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = couleur_accent
    shape.line.fill.background()

    ajouter_texte(
        slide, titre,
        0.6, 0.2, 9, 1,
        26, couleur_accent, bold=True
    )

    # Ligne sous le titre
    shape2 = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.15), Inches(8.8), Inches(0.03)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = FOND_CLAIR
    shape2.line.fill.background()

    if contenu and contenu.strip():
        ajouter_texte(
            slide, contenu,
            0.6, 1.3, 9, 1.2,
            15, GRIS, wrap=True
        )

    if points:
        y = 2.6 if contenu else 1.5
        for point in points[:5]:
            if point and str(point).strip():
                ajouter_texte(
                    slide, "▶  " + str(point),
                    0.8, y, 8.8, 0.75,
                    14, BLANC, wrap=True
                )
                y += 0.82

def slide_resultats(prs, titre, contenu, points=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, FOND_SOMBRE)

    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = VERT
    shape.line.fill.background()

    ajouter_texte(
        slide, titre,
        0.6, 0.2, 9, 1,
        26, VERT, bold=True
    )

    if contenu:
        ajouter_texte(
            slide, contenu,
            0.6, 1.3, 9, 1.2,
            15, GRIS, wrap=True
        )

    if points:
        y = 2.6
        for point in points[:5]:
            if point and str(point).strip():
                ajouter_texte(
                    slide, "✅  " + str(point),
                    0.8, y, 8.8, 0.75,
                    14, BLANC, wrap=True
                )
                y += 0.82

def slide_conclusion(prs, titre, contenu, points=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, VIOLET)

    ajouter_texte(
        slide, titre,
        1, 1.2, 8, 1.2,
        32, BLANC, bold=True,
        align=PP_ALIGN.CENTER
    )

    if contenu:
        ajouter_texte(
            slide, contenu,
            1, 2.6, 8, 1.5,
            16, BLANC,
            align=PP_ALIGN.CENTER, wrap=True
        )

    if points:
        y = 4.3
        for point in points[:3]:
            if point and str(point).strip():
                ajouter_texte(
                    slide, "▶  " + str(point),
                    1, y, 8, 0.7,
                    14, JAUNE, wrap=True
                )
                y += 0.75

def slide_remerciements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fond(slide, FOND_SOMBRE)

    ajouter_texte(
        slide, "Merci pour votre attention ! 🙏",
        1, 2.5, 8, 1.5,
        36, ORANGE, bold=True,
        align=PP_ALIGN.CENTER
    )
    ajouter_texte(
        slide, "Des questions ?",
        1, 4, 8, 1,
        22, BLANC,
        align=PP_ALIGN.CENTER
    )
    ajouter_texte(
        slide, "Généré par Mpianatra-midev",
        1, 6.5, 8, 0.6,
        11, GRIS,
        align=PP_ALIGN.CENTER
    )

def create_pptx(text: str, output_path: str,
                slides_data: list = None,
                analyse: dict = None):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    if slides_data and len(slides_data) > 0:

        # 1. Slide salutation
        slide_salutation(prs)

        # 2. Slide titre
        titre_slide = next(
            (s for s in slides_data if s.get("type") == "titre"),
            slides_data[0]
        )
        slide_titre(
            prs,
            titre_slide.get("titre", "Présentation"),
            titre_slide.get("auteur", ""),
            titre_slide.get("institution", ""),
            titre_slide.get("annee", "")
        )

        # 3. Slide sommaire
        slide_sommaire(prs, slides_data)

        # 4. Slides contenu
        couleurs = [ORANGE, RGBColor(56, 189, 248),
                    RGBColor(167, 139, 250), JAUNE,
                    RGBColor(251, 113, 133)]
        couleur_idx = 0

        for slide in slides_data:
            type_slide = str(slide.get("type", "contenu")).lower()
            titre = str(slide.get("titre", ""))
            contenu = str(slide.get("contenu", ""))
            points = slide.get("points", None)

            if type_slide in ["titre", "sommaire"]:
                continue
            elif type_slide == "resultats":
                slide_resultats(prs, titre, contenu, points)
            elif type_slide in ["conclusion", "discussion"]:
                slide_conclusion(prs, titre, contenu, points)
            elif type_slide == "remerciements":
                slide_remerciements(prs)
            else:
                couleur = couleurs[couleur_idx % len(couleurs)]
                slide_contenu(prs, titre, contenu, points, couleur)
                couleur_idx += 1

        # 5. Slide remerciements finale
        slide_remerciements(prs)

    else:
        # Fallback sans IA
        slide_salutation(prs)
        lines = [l.strip() for l in text.split("\n") if len(l.strip()) > 20]
        slide_titre(prs, "Présentation", "", "", "")
        for i, line in enumerate(lines[:10]):
            slide_contenu(prs, f"Point {i+1}", line)
        slide_remerciements(prs)

    prs.save(output_path)
    print(f"✅ PowerPoint sauvegardé : {output_path}")
    return output_path